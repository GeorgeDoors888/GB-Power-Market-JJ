from src.extraction.docx_parser import extract_docx_text
from src.extraction.pptx_parser import extract_pptx_text
from docx import Document
from pptx import Presentation
import io

def make_docx_bytes():
    buf = io.BytesIO()
    doc = Document()
    doc.add_paragraph("Hello world from DOCX")
    doc.save(buf)
    return buf.getvalue()

def make_pptx_bytes():
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Hello PPTX"
    prs.save(buf)
    return buf.getvalue()

def test_extract_docx_text():
    b = make_docx_bytes()
    txt = extract_docx_text(b)
    assert "Hello world" in txt

def test_extract_pptx_text():
    b = make_pptx_bytes()
    txt = extract_pptx_text(b)
    assert "Hello PPTX" in txt

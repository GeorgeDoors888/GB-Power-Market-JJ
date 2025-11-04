from __future__ import annotations
import io
from pptx import Presentation

def extract_pptx_text(blob: bytes) -> str:
    prs = Presentation(io.BytesIO(blob))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)
